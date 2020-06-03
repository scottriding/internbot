import pptx
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.fill import FillFormat
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from collections import OrderedDict

class Powerpoint(object):

    def pick_template_layout(self, path_to_template):
        self.__presentation = pptx.Presentation(path_to_template)
        template_names = []
        for layout in self.__presentation.slide_layouts:
            template_names.append(layout.name)
        return template_names

    def build_powerpoint_report(self, question_blocks, layout_index, path_to_output):
        self.__question_blocks = question_blocks
        self.__layout_index = layout_index

        self.write_questions()
        print("Saving report...")
        self.__presentation.save(path_to_output)

    def write_questions(self):
        for question in self.__question_blocks.questions:
            to_print = "Writing question: %s" % question.name
            print(to_print)

            if question.parent == "CompositeQuestion":
                if question.type == "CompositeMatrix":
                    self.write_matrix_question(question)
                else:
                    self.write_binary_question(question)
            elif question.type == 'TE':
                pass
            elif question.type == "Slider":
                self.write_basic_charts(question)
            else:
                self.write_question(question)

    def write_matrix_question(self, matrix_question):
        ## if any of the sub_questions are grouped, break them all up
        break_up = False
        for sub_question in matrix_question.questions:
            total_ns = self.pull_ns(sub_question)
            if len(total_ns) > 1:
                break_up = True

        if break_up:
            for sub_question in matrix_question.questions:
                self.write_question(sub_question)
        else:
            self.write_combined_matrix_question(matrix_question)

    def write_combined_matrix_question(self, question):
        matrix_slide_one = self.__presentation.slides.add_slide(self.__presentation.slide_layouts[self.__layout_index])
        matrix_title = matrix_slide_one.shapes.title
        if matrix_title:
            matrix_title.text = question.name

        matrix_slide_two = self.__presentation.slides.add_slide(self.__presentation.slide_layouts[self.__layout_index])
        matrix_title = matrix_slide_two.shapes.title
        if matrix_title:
            matrix_title.text = question.name

        matrix_slide_three = self.__presentation.slides.add_slide(self.__presentation.slide_layouts[self.__layout_index])
        matrix_title = matrix_slide_three.shapes.title
        if matrix_title:
            matrix_title.text = question.name

        series_names = []  # scale points
        categories = []  # statements for matrix
        series_lists = OrderedDict()  # frequencies for each scale point

        for response in question.questions[0].responses:
            series_names.append(response.label)

        for sub_question in question.questions:
            categories.append(sub_question.prompt)
            for response in sub_question.responses:
                if series_lists.get(response.label) is not None:
                    if len(response.frequencies.frequencies.items()) > 0:
                        for group, frequency in response.frequencies.frequencies.items():
                            if frequency.result == "NA":
                                series_lists[response.label].append(0)
                            else:
                                series_lists[response.label].append(frequency.result)
                    else:
                        series_lists[response.label].append(0)
                else:
                    if len(response.frequencies.frequencies.items()) > 0:
                        for group, frequency in response.frequencies.frequencies.items():
                            if frequency.result == "NA":
                                series_lists[response.label] = [0]
                            else:
                                series_lists[response.label] = [frequency.result]
                    else:
                        series_lists[response.label].append(0)

        self.column_clustered_matrix_by_statements(question, matrix_slide_one, categories, series_names, series_lists)
        self.bar_clustered_matrix_by_statements(question, matrix_slide_two, categories, series_names, series_lists)
        self.stacked_bar_matrix(question, matrix_slide_three, categories, series_names, series_lists, True)
        self.stacked_bar_matrix(question, matrix_slide_three, categories, series_names, series_lists, False)

        series_names = []  # statements
        categories = []  # scale points
        series_lists = OrderedDict()  # frequencies for each scale point

        for response in question.questions[0].responses:
            series_names.append(response.label)

        for sub_question in question.questions:
            categories.append(sub_question.prompt)
            if series_lists.get(sub_question.prompt) is not None:
                current = series_lists.get(sub_question.prompt)
            else:
                current = []
            for response in sub_question.responses:
                if len(response.frequencies.frequencies.items()) > 0:
                    for group, frequency in response.frequencies.frequencies.items():
                        if frequency.result == "NA":
                            current.append(0)
                        else:
                            current.append(frequency.result)
                else:
                    current.append(0)
            series_lists[sub_question.prompt] = current

        self.column_clustered_matrix_by_scale_points(question, matrix_slide_one, categories, series_names, series_lists)
        self.bar_clustered_matrix_by_scale_points(question, matrix_slide_two, categories, series_names, series_lists)

    def column_clustered_matrix_by_statements(self, question, slide, categories, series_names, series_lists):
        chart_data = CategoryChartData()
        chart_data.categories = categories
        
        for name in series_names:
            chart_data.add_series(name, iter(series_lists[name]))

        left, top, width, height = Inches(0), Inches(2), Inches(6.5), Inches(5)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data).chart
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = "Q: " + str(question.prompt)
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)

        category_axis = chart.category_axis
        category_axis.has_minor_gridlines = False
        category_axis.has_major_gridlines = False
        category_axis.tick_labels.font.size = Pt(12)
        value_axis = chart.value_axis
        value_axis.has_minor_gridlines = False
        value_axis.has_major_gridlines = False
        value_axis.tick_labels.font.size = Pt(12)
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        data_labels.font.size = Pt(12)
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(12)
        chart.legend.include_in_layout = False

    def column_clustered_matrix_by_scale_points(self, question, slide, categories, series_names, series_lists):
        chart_data = CategoryChartData()
        chart_data.categories = series_names
        for category in categories:
            chart_data.add_series(category, iter(series_lists[category]))

        left, top, width, height = Inches(6.75), Inches(2), Inches(6.5), Inches(5)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data).chart
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = "Q: " + str(question.prompt)
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)

        category_axis = chart.category_axis
        category_axis.has_minor_gridlines = False
        category_axis.has_major_gridlines = False
        category_axis.tick_labels.font.size = Pt(12)
        value_axis = chart.value_axis
        value_axis.has_minor_gridlines = False
        value_axis.has_major_gridlines = False
        value_axis.tick_labels.font.size = Pt(12)
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        data_labels.font.size = Pt(12)
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(12)
        chart.legend.include_in_layout = False

    def bar_clustered_matrix_by_statements(self, question, slide, categories, series_names, series_lists):
        chart_data = CategoryChartData()
        categories.reverse()
        chart_data.categories = categories

        series_names.reverse()
        for name in series_names:
            new_list = series_lists[name]
            new_list.reverse()
            chart_data.add_series(name, iter(new_list))

        left, top, width, height = Inches(0), Inches(2), Inches(6.5), Inches(5)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data).chart
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = "Q: " + str(question.prompt)
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)

        category_axis = chart.category_axis
        category_axis.has_minor_gridlines = False
        category_axis.has_major_gridlines = False
        category_axis.tick_labels.font.size = Pt(12)
        value_axis = chart.value_axis
        value_axis.has_minor_gridlines = False
        value_axis.has_major_gridlines = False
        value_axis.tick_labels.font.size = Pt(12)
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        data_labels.font.size = Pt(12)
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(12)
        chart.legend.include_in_layout = False

    def bar_clustered_matrix_by_scale_points(self, question, slide, categories, series_names, series_lists):
        series_names.reverse()
        chart_data = CategoryChartData()
        chart_data.categories = series_names

        categories.reverse()
        for category in categories:
            new_list = series_lists[category]
            new_list.reverse()
            chart_data.add_series(category, iter(new_list))

        left, top, width, height = Inches(6.75), Inches(2), Inches(6.5), Inches(5)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data).chart
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = "Q: " + str(question.prompt)
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)

        category_axis = chart.category_axis
        category_axis.has_minor_gridlines = False
        category_axis.has_major_gridlines = False
        category_axis.tick_labels.font.size = Pt(12)
        value_axis = chart.value_axis
        value_axis.has_minor_gridlines = False
        value_axis.has_major_gridlines = False
        value_axis.tick_labels.font.size = Pt(12)
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        data_labels.font.size = Pt(12)
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(12)
        chart.legend.include_in_layout = False

    def stacked_bar_matrix(self, question, slide, categories, series_names, series_lists, is_right):
        chart_data = CategoryChartData()
        chart_data.categories = categories

        for name in series_names:
            chart_data.add_series(name, iter(series_lists[name]))

        if is_right:
            left, top, width, height = Inches(0), Inches(2), Inches(6.5), Inches(5)
        else:
            left, top, width, height = Inches(6.75), Inches(2), Inches(6.5), Inches(5)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED, left, top, width, height, chart_data).chart
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = "Q: " + str(question.prompt)
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
        chart.has_legend = True
        category_axis = chart.category_axis
        category_axis.has_minor_gridlines = False
        category_axis.has_major_gridlines = False
        category_axis.tick_labels.font.size = Pt(12)
        value_axis = chart.value_axis
        value_axis.has_minor_gridlines = False
        value_axis.has_major_gridlines = False
        value_axis.visible = False

        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(12)
        chart.legend.include_in_layout = False
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.number_format = '0%'
        data_labels.font.size = Pt(12)
        plot = chart.plots[0]
        num_series = len(plot.series)


        if num_series <= 6:
            colors = []
            if num_series >= 1:
                colors.append(1)
            if num_series >= 2:
                colors.append(6)
            if num_series >= 3:
                colors.insert(1, 4)
            if num_series >= 4:
                colors.insert(1, 2)
            if num_series >= 5:
                colors.insert(3, 5)
            if num_series == 6:
                colors.insert(2, 3)

            if is_right:
                colors.reverse()

            color_index = 0
            for series in plot.series:
                fill = series.format.fill
                fill.solid()
                if colors[color_index] == 1:
                    fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
                elif colors[color_index] == 2:
                    fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_2
                elif colors[color_index] == 3:
                    fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_3
                elif colors[color_index] == 4:
                    fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_4
                elif colors[color_index] == 5:
                    fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_5
                elif colors[color_index] == 6:
                    fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_6
                color_index += 1
        else:
            print("CANNOT REVERSE GRADIENT FOR STACKED BARS WITH MORE THAN 6 SERIES")

    def write_binary_question(self, binary_question):
        groups = self.pull_binary_groups(binary_question)
        if len(groups) == 1:
            pass  
        else:
            pass

    def write_question(self, question):
        total_ns = self.pull_ns(question)
        if len(total_ns) == 1:
            if len(question.responses) < 3:
                self.write_pie_chart(question)
            self.write_basic_charts(question)
        else:
            self.write_grouped_charts(question, total_ns)

    def write_pie_chart(self, question):
        pie_slide = self.__presentation.slides.add_slide(self.__presentation.slide_layouts[1])
        title = pie_slide.shapes.title
        if title:
            title.text = question.name

        categories = []
        frequencies = []
        is_percent = False
        for response in question.responses:
            categories.append(str(response.label))
            for group, frequency in response.frequencies.frequencies.items():
                if frequency.stat == 'percent':
                    is_percent = True
                if frequency.result == "NA":
                    frequencies.append(0)
                else:
                    frequencies.append(frequency.result)

        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series('Series 1', iter(frequencies))
        left, top, width, height = Inches(6.75), Inches(2), Inches(6.5), Inches(5)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, left, top, width, height, chart_data).chart
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = "Q: " + str(question.prompt)
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(12)
        chart.legend.include_in_layout = False
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        if is_percent:
            data_labels.number_format = '0%'
        data_labels.font.size = Pt(12)

    def write_basic_charts(self, question):
        bar_slide_one = self.__presentation.slides.add_slide(self.__presentation.slide_layouts[self.__layout_index])
        bar_title = bar_slide_one.shapes.title
        if bar_title:
            bar_title.text = question.name

        bar_slide_two = self.__presentation.slides.add_slide(self.__presentation.slide_layouts[self.__layout_index])
        bar_title = bar_slide_two.shapes.title
        if bar_title:
            bar_title.text = question.name

        bar_slide_three = self.__presentation.slides.add_slide(self.__presentation.slide_layouts[self.__layout_index])
        bar_title = bar_slide_three.shapes.title
        if bar_title:
            bar_title.text = question.name

        bar_slide_four = self.__presentation.slides.add_slide(self.__presentation.slide_layouts[self.__layout_index])
        bar_title = bar_slide_four.shapes.title
        if bar_title:
            bar_title.text = question.name

        bar_slide_five = self.__presentation.slides.add_slide(self.__presentation.slide_layouts[self.__layout_index])
        bar_title = bar_slide_five.shapes.title
        if bar_title:
            bar_title.text = question.name

        bar_slide_six = self.__presentation.slides.add_slide(self.__presentation.slide_layouts[self.__layout_index])
        bar_title = bar_slide_six.shapes.title
        if bar_title:
            bar_title.text = question.name

        stacked_slide = self.__presentation.slides.add_slide(self.__presentation.slide_layouts[self.__layout_index])
        stacked_title = stacked_slide.shapes.title
        if stacked_title:
            stacked_title.text = question.name

        categories = []
        frequencies = []
        is_percent = False
        for response in question.responses:
            categories.append(str(response.label))
            for group, frequency in response.frequencies.frequencies.items():
                if frequency.stat == 'percent':
                    is_percent = True
                if frequency.result == "NA":
                    frequencies.append(0)
                else:
                    frequencies.append(frequency.result)

        ## all vertical charts in all colors
        self.vertical_bars(bar_slide_one, question.prompt, categories, frequencies, MSO_THEME_COLOR.ACCENT_1, is_percent)
        self.vertical_bars(bar_slide_two, question.prompt, categories, frequencies, MSO_THEME_COLOR.ACCENT_2, is_percent)
        self.vertical_bars(bar_slide_three, question.prompt, categories, frequencies, MSO_THEME_COLOR.ACCENT_3, is_percent)
        self.vertical_bars(bar_slide_four, question.prompt, categories, frequencies, MSO_THEME_COLOR.ACCENT_4, is_percent)
        self.vertical_bars(bar_slide_five, question.prompt, categories, frequencies, MSO_THEME_COLOR.ACCENT_5, is_percent)
        self.vertical_bars(bar_slide_six, question.prompt, categories, frequencies, MSO_THEME_COLOR.ACCENT_6, is_percent)


        ## all stacked charts in all colors
        self.stacked_bars(stacked_slide, question.prompt, categories, frequencies, is_percent)

        categories.reverse()
        frequencies.reverse()

        ## all horizontal charts in all colors
        self.horizontal_bars(bar_slide_one, question.prompt, categories, frequencies, MSO_THEME_COLOR.ACCENT_1, is_percent)
        self.horizontal_bars(bar_slide_two, question.prompt, categories, frequencies, MSO_THEME_COLOR.ACCENT_2, is_percent)
        self.horizontal_bars(bar_slide_three, question.prompt, categories, frequencies, MSO_THEME_COLOR.ACCENT_3, is_percent)
        self.horizontal_bars(bar_slide_four, question.prompt, categories, frequencies, MSO_THEME_COLOR.ACCENT_4, is_percent)
        self.horizontal_bars(bar_slide_five, question.prompt, categories, frequencies, MSO_THEME_COLOR.ACCENT_5, is_percent)
        self.horizontal_bars(bar_slide_six, question.prompt, categories, frequencies, MSO_THEME_COLOR.ACCENT_6, is_percent)

    def vertical_bars(self, slide, prompt, categories, frequencies, color, is_percent):
        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series('Series 1', iter(frequencies))

        left, top, width, height = Inches(0), Inches(2), Inches(6.5), Inches(5)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data).chart
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = "Q: " + str(prompt)
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
        chart.has_legend = False
        category_axis = chart.category_axis
        category_axis.has_minor_gridlines = False
        category_axis.has_major_gridlines = False
        category_axis.tick_labels.font.size = Pt(12)
        value_axis = chart.value_axis
        value_axis.has_minor_gridlines = False
        value_axis.has_major_gridlines = False
        value_axis.tick_labels.font.size = Pt(12)
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        if is_percent:
            data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        data_labels.font.size = Pt(12)
        plot = chart.plots[0]
        series = plot.series[0]
        fill = series.format.fill
        fill.solid()
        fill.fore_color.theme_color = color

    def stacked_bars(self, slide, prompt, categories, frequencies, is_percent):
        chart_data = CategoryChartData()
        chart_data.categories = ["Category 1"]

        if len(frequencies) > 0:
            for index in range(0, len(categories)):
                temp = []
                temp.append(frequencies[index])
                series = str(categories[index])
                chart_data.add_series(series, iter(temp))
        else:
            temp = [0]
            chart_data.add_series('Category 1', iter(temp))

        left, top, width, height = Inches(0), Inches(2), Inches(6.5), Inches(5)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED, left, top, width, height, chart_data).chart
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = "Q: " + str(prompt)
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
        chart.has_legend = True
        category_axis = chart.category_axis
        category_axis.has_minor_gridlines = False
        category_axis.has_major_gridlines = False
        category_axis.tick_labels.font.size = Pt(12)
        value_axis = chart.value_axis
        value_axis.has_minor_gridlines = False
        value_axis.has_major_gridlines = False
        value_axis.visible = False

        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(12)
        chart.legend.include_in_layout = False
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        if is_percent:
            data_labels.number_format = '0%'
        data_labels.font.size = Pt(12)

    def horizontal_bars(self, slide, prompt, categories, frequencies, color, is_percent):
        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series('Series 1', iter(frequencies))

        left, top, width, height = Inches(6.75), Inches(2), Inches(6.5), Inches(5)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data).chart
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = "Q: " + str(prompt)
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
        chart.has_legend = False
        category_axis = chart.category_axis
        category_axis.has_minor_gridlines = False
        category_axis.has_major_gridlines = False
        category_axis.tick_labels.font.size = Pt(12)
        value_axis = chart.value_axis
        value_axis.has_minor_gridlines = False
        value_axis.has_major_gridlines = False
        value_axis.tick_labels.font.size = Pt(12)
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        if is_percent:
            data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        data_labels.font.size = Pt(12)
        plot = chart.plots[0]
        series = plot.series[0]
        fill = series.format.fill
        fill.solid()
        fill.fore_color.theme_color = color

    def write_grouped_charts(self, question, total_ns):
        grouped_bars_one = self.__presentation.slides.add_slide(self.__presentation.slide_layouts[self.__layout_index])
        title = grouped_bars_one.shapes.title
        if title:
            title.text = question.name

        grouped_bars_two = self.__presentation.slides.add_slide(self.__presentation.slide_layouts[self.__layout_index])
        title = grouped_bars_two.shapes.title
        if title:
            title.text = question.name

        grouped_bars_three = self.__presentation.slides.add_slide(self.__presentation.slide_layouts[self.__layout_index])
        title = grouped_bars_three.shapes.title
        if title:
            title.text = question.name

        series_names = list(total_ns.keys())  # groups
        categories = []  # responses
        series_lists = OrderedDict()  # frequencies for each group

        for response in question.responses:
            categories.append(response.label)
            for name in series_names:
                if series_lists.get(name) is not None:
                    current = series_lists.get(name)
                else:
                    current = []
                if len(response.frequencies.frequencies.items()) > 0:
                    if response.frequencies.frequencies.get(name) is not None:
                        frequency = response.frequencies.frequencies.get(name)
                        if frequency.result == "NA":
                            current.append(0)
                        else:
                            current.append(frequency.result)
                    else:
                        current.append(0)
                else:
                    current.append(0)
                series_lists[name] = current

        self.column_clustered_matrix_by_statements(question, grouped_bars_one, categories, series_names, series_lists)
        self.bar_clustered_matrix_by_statements(question, grouped_bars_two, categories, series_names, series_lists)
        self.stacked_bar_matrix(question, grouped_bars_three, categories, series_names, series_lists, True)
        self.stacked_bar_matrix(question, grouped_bars_three, categories, series_names, series_lists, False)

        series_names = []  # responses
        categories = list(total_ns.keys())  # groups
        series_lists = OrderedDict()  # frequencies for each response

        for response in question.responses:
            series_names.append(response.label)
            for name in categories:
                if series_lists.get(name) is not None:
                    current = series_lists.get(name)
                else:
                    current = []
                if len(response.frequencies.frequencies.items()) > 0:
                    if response.frequencies.frequencies.get(name) is not None:
                        frequency = response.frequencies.frequencies.get(name)
                        if frequency.result == "NA":
                            current.append(0)
                        else:
                            current.append(frequency.result)
                    else:
                        current.append(0)
                else:
                    current.append(0)
                series_lists[name] = current

        self.column_clustered_matrix_by_scale_points(question, grouped_bars_one, categories, series_names, series_lists)
        self.bar_clustered_matrix_by_scale_points(question, grouped_bars_two, categories, series_names, series_lists)

    def pull_ns(self, question):
        ## we're going to total all the groups that are in the responses of this question
        total_ns = OrderedDict()
        for response in question.responses:
            for group, frequency in response.frequencies.frequencies.items():
                if frequency.result != "NA":
                ### if the response belongs to a group, calculate the n
                    if total_ns.get(group) is not None:
                        total_ns[group] += int(frequency.population)
                    else:
                        total_ns[group] = int(frequency.population)
        return total_ns

    def pull_binary_groups(self, binary_question):
        ## we're going to total all the groups that are in the responses of this question
        total_groups = []
        for sub_question in binary_question.questions:
            for response in sub_question.responses:
                for group, frequency in response.frequencies.frequencies.items():
                    if frequency.result != "NA":
                        if group not in total_groups:
                            total_groups.append(group)
        return total_groups

    
