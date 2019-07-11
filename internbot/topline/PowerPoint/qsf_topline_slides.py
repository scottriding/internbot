
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
class QSFToplineSlides(object):

    def __init__ (self, survey, path_to_template, path_to_output, years):
        self.presentation = Presentation(path_to_template)
        self.questions = survey.get_questions()
        self.survey = survey
        self.years = years
        self.path_to_output = path_to_output
        self.slide_layouts = self.presentation.slide_layouts
        self.diverge_color_array = [RGBColor(150, 215, 5),
                                    RGBColor(190, 230, 100),
                                    RGBColor(210, 240, 150),
                                    RGBColor(230, 245, 190),
                                    RGBColor(235, 235, 240),
                                    RGBColor(215, 220, 235),
                                    RGBColor(190, 205, 230),
                                    RGBColor(170, 185, 245),
                                    RGBColor(135, 145, 190),
                                    RGBColor(0, 35, 100)]

        self.mono_color_array = [RGBColor(150, 215, 5),
                                 RGBColor(160, 125, 50),
                                 RGBColor(190, 230, 100),
                                 RGBColor(210, 240, 150),
                                 RGBColor(230, 245, 190),
                                 RGBColor(235, 235, 240),
                                 RGBColor(215, 215, 220),
                                 RGBColor(190, 190, 195),
                                 RGBColor(165, 165, 170),
                                 RGBColor(130, 130, 140)]

        self.colors_to_use = []

    def save(self, path_to_output):
        self.chart_questions()
        self.presentation.save(path_to_output)
        print("Finished!")

    def chart_questions(self):
        title_slide = self.presentation.slides[0]
        title_slide.shapes[0].text = self.survey.name
        title_slide.shapes[1].text = "Y2 Analytics Slide Automation"
        for question in self.questions:
            to_print = "Writing question: %s" % question.name
            print(to_print)
            try:
                if question.parent == 'CompositeQuestion':
                    self.chart_composite_question(question)
                elif question.type == 'TE':
                    pass
                else:
                    self.chart_question(question)
            except ValueError:
                pass

    def chart_question(self, question):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts.get_by_name('AutomatedChart'))
        shapes = slide.shapes
        self.write_question_details(slide, question)

        if len(self.years) > 0:
            self.column_clustered_basic_trended(question, slide)
            self.bar_clustered_basic_trended(question, slide)
        else:
            if len(question.responses) < 3:
                self.pie_simple(question, slide)
                self.column_clustered_basic(question, slide)
            else:
                self.column_clustered_basic(question, slide)
                self.bar_clustered_basic(question, slide)

    def chart_composite_question(self, question):
        if len(self.years) > 0:
            pass
        else:
            if question.type == 'CompositeMatrix':
                slide_1 = self.presentation.slides.add_slide(self.presentation.slide_layouts.get_by_name('AutomatedChart'))
                self.write_composite_question_details(slide_1, question, 1, 4)
                self.column_clustered_matrix_by_categories(question, slide_1)
                self.column_clustered_matrix_by_scale_points(question, slide_1)

                slide_2 = self.presentation.slides.add_slide(self.presentation.slide_layouts.get_by_name('AutomatedChart'))
                self.write_composite_question_details(slide_2, question, 2, 4)
                self.bar_clustered_matrix_by_categories(question, slide_2)
                self.bar_clustered_matrix_by_scale_points(question, slide_2)

                slide_3 = self.presentation.slides.add_slide(self.presentation.slide_layouts.get_by_name('AutomatedChart'))
                self.write_composite_question_details(slide_3, question, 3, 4)
                self.stacked_bar_matrix(question, slide_3, True, True, 'top')
                self.stacked_bar_matrix(question, slide_3, False, True, 'bottom')

                slide_4 = self.presentation.slides.add_slide(self.presentation.slide_layouts.get_by_name('AutomatedChart'))
                self.write_composite_question_details(slide_4, question, 4, 4)
                self.stacked_bar_matrix(question, slide_4, True, False, 'top')
                self.stacked_bar_matrix(question, slide_4, False, False, 'bottom')




            elif question.type == 'CompositeConstantSum':
                slide_1 = self.presentation.slides.add_slide(self.presentation.slide_layouts.get_by_name('AutomatedChart'))
                self.write_composite_question_details(slide_1, question, 1, 1)
                self.column_clustered_allocation(question, slide_1)
                self.bar_clustered_allocation(question, slide_1)
            else:
                # self.write_binary(question.questions)
                pass

    def write_question_details(self, slide, question):
        header = slide.shapes[0]
        header.text = str(question.name)
        n_box = slide.shapes[1]
        n_box.text = "N=" + str(question.n)
        count_box = slide.shapes[2]
        count_box.text = "Slide 1 of 1"

    def write_composite_question_details(self, slide, question, slide_number, slide_total):
        header = slide.shapes[0]
        header.text = str(question.name)
        n_box = slide.shapes[1]
        n_box.text = "n's in charts"
        count_box = slide.shapes[2]
        count_box.text = "Slide "+str(slide_number)+" of "+str(slide_total)

    def max_years(self, responses):
        years_used = []
        for response in responses:
            for year in self.years:
                if response.frequencies.get(year) is not None:
                    if year not in years_used:
                        years_used.append(year)
        return years_used

    def max_years_subquestions(self, sub_questions):
        years_used = []
        for question in sub_questions:
            for response in question.responses:
                for year in self.years:
                    if response.frequencies.get(year) is not None:
                        if year not in years_used:
                            years_used.append(year)
        return years_used


    def column_clustered_basic(self, question, slide):
        categories = []
        frequencies = []
        for response in question.responses:
            categories.append(str(response.response))
            for year, frequency in response.frequencies.items():
                frequencies.append(frequency)
        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series('Series 1', iter(frequencies))
        left, top, width, height = Inches(0), Inches(2), Inches(6.5), Inches(5)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data).chart
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = "Q: " + str(question.prompt)
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
        chart.has_legend = False
        category_axis = chart.category_axis
        category_axis.has_minor_gridlines = False
        category_axis.has_major_gridlines = False
        category_axis.tick_labels.font.size = Pt(12)
        value_axis = chart.value_axis
        value_axis.has_minor_gridlines = False
        value_axis.has_major_gridlines = False
        value_axis.tick_labels.font.size = Pt(12)
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        if question.stat == 'percent':
            data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        data_labels.font.size = Pt(12)
        plot = chart.plots[0]
        series = plot.series[0]
        fill = series.format.fill
        fill.solid()
        fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_3

    def column_clustered_basic_trended(self, question, slide):
        headers = self.max_years(question.responses)
        series_lists = []
        categories = []
        for header in headers:
            frequencies = []
            series_lists.append(frequencies)
        for response in question.responses:
            categories.append(str(response.response))
            header_idx = 0
            for header in headers:
                frequencies = series_lists[header_idx]
                if response.frequencies.get(header) is not None:
                    freq = response.frequencies.get(header)
                    frequencies.append(freq)
                header_idx += 1
        chart_data = CategoryChartData()
        chart_data.categories = categories
        for idx in range(0, len(headers)):
            chart_data.add_series(str(headers[idx]), iter(series_lists[idx]))

        left, top, width, height = Inches(0), Inches(2), Inches(6.5), Inches(5)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data).chart
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = "Q: " + str(question.prompt)
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
        chart.has_legend = False
        category_axis = chart.category_axis
        category_axis.has_minor_gridlines = False
        category_axis.has_major_gridlines = False
        category_axis.tick_labels.font.size = Pt(12)
        value_axis = chart.value_axis
        value_axis.has_minor_gridlines = False
        value_axis.has_major_gridlines = False
        value_axis.tick_labels.font.size = Pt(12)
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        if question.stat == 'percent':
            data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        data_labels.font.size = Pt(12)
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(12)
        chart.legend.include_in_layout = False

    def bar_clustered_basic(self, question, slide):
        categories = []
        frequencies = []
        for response in question.responses:
            categories.append(str(response.response))
            for year, frequency in response.frequencies.items():
                frequencies.append(frequency)
        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series('Series 1', iter(frequencies))
        left, top, width, height = Inches(6.75), Inches(2), Inches(6.5), Inches(5)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data).chart
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = "Q: " + str(question.prompt)
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
        chart.has_legend = False
        category_axis = chart.category_axis
        category_axis.has_minor_gridlines = False
        category_axis.has_major_gridlines = False
        category_axis.tick_labels.font.size = Pt(12)
        value_axis = chart.value_axis
        value_axis.has_minor_gridlines = False
        value_axis.has_major_gridlines = False
        value_axis.tick_labels.font.size = Pt(12)
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        if question.stat == 'percent':
            data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        data_labels.font.size = Pt(12)
        plot = chart.plots[0]
        series = plot.series[0]
        fill = series.format.fill
        fill.solid()
        fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_2

    def bar_clustered_basic_trended(self, question, slide):
        headers = self.max_years(question.responses)
        series_lists = []
        categories = []
        for header in headers:
            frequencies = []
            series_lists.append(frequencies)
        for response in question.responses:
            categories.append(str(response.response))
            header_idx = 0
            for header in headers:
                frequencies = series_lists[header_idx]
                if response.frequencies.get(header) is not None:
                    freq = response.frequencies.get(header)
                    frequencies.append(freq)
                header_idx += 1
        chart_data = CategoryChartData()
        chart_data.categories = categories
        for idx in range(0, len(headers)):
            chart_data.add_series(str(headers[idx]), iter(series_lists[idx]))
        left, top, width, height = Inches(6.75), Inches(2), Inches(6.5), Inches(5)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data).chart
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = "Q: " + str(question.prompt)
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
        chart.has_legend = False
        category_axis = chart.category_axis
        category_axis.has_minor_gridlines = False
        category_axis.has_major_gridlines = False
        category_axis.tick_labels.font.size = Pt(12)
        value_axis = chart.value_axis
        value_axis.has_minor_gridlines = False
        value_axis.has_major_gridlines = False
        value_axis.tick_labels.font.size = Pt(12)
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        if question.stat == 'percent':
            data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        data_labels.font.size = Pt(12)
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(12)
        chart.legend.include_in_layout = False

    def stacked_bar_simple(self, question, slide):
        series_names = []
        frequencies = []
        for response in question.responses:
            series_names.append(str(response.response))
            for year, frequency in response.frequencies.items():
                frequencies.append(frequency)
        chart_data = CategoryChartData()
        chart_data.categories = ["Variable 1"]
        for idx in range(0, len(series_names)):
            temp = []
            temp.append(frequencies[idx])
            series = str(series_names[idx])
            chart_data.add_series(series, iter(temp))

        left, top, width, height = Inches(0), Inches(4), Inches(6.5), Inches(3)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED, left, top, width, height, chart_data).chart
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = "Q: " + str(question.prompt)
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
        chart.has_legend = True
        category_axis = chart.category_axis
        category_axis.has_minor_gridlines = False
        category_axis.has_major_gridlines = False
        category_axis.tick_labels.font.size = Pt(12)
        value_axis = chart.value_axis
        value_axis.has_minor_gridlines = False
        value_axis.has_major_gridlines = False
        value_axis.visible = False

        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(12)
        chart.legend.include_in_layout = False
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        if question.stat == 'percent':
            data_labels.number_format = '0%'
        data_labels.font.size = Pt(12)

    def pie_simple(self, question, slide):
        categories = []
        frequencies = []
        for response in question.responses:
            categories.append(str(response.response))
            for year, frequency in response.frequencies.items():
                frequencies.append(frequency)
        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series('Series 1', iter(frequencies))
        left, top, width, height = Inches(6.75), Inches(2), Inches(6.5), Inches(5)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, left, top, width, height, chart_data).chart
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = "Q: " + str(question.prompt)
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(12)
        chart.legend.include_in_layout = False
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        if question.stat == 'percent':
            data_labels.number_format = '0%'
        data_labels.font.size = Pt(12)

    def column_clustered_matrix_by_categories(self, question, slide):
        sub_questions = question.questions
        series_names = []  # scale points
        series_lists = []  # frequencies for each scale point
        categories = []  # responses for matrix
        is_first = True

        for sub_question in sub_questions:
            categories.append(str(sub_question.prompt) + " (n=" + str(sub_question.n) + ")")
            if is_first:
                for response in sub_question.responses:
                    series_names.append(response.response)
                question.stat = sub_question.stat
            is_first = False

        for name in series_names:
            frequencies = []
            series_lists.append([])

        for sub_question in sub_questions:
            response_idx = 0
            for response in sub_question.responses:
                frequencies = series_lists[response_idx]
                if response.has_frequency is True:
                    for year, frequency in response.frequencies.items():
                        frequencies.append(frequency)
                else:
                    frequencies.append(0)
                response_idx += 1

        chart_data = CategoryChartData()
        chart_data.categories = categories
        for idx in range(0, len(series_names)):
            chart_data.add_series(series_names[idx], iter(series_lists[idx]))

        left, top, width, height = Inches(0), Inches(2), Inches(6.5), Inches(5)
        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data).chart
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = "Q: " + str(question.prompt)
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)

        category_axis = chart.category_axis
        category_axis.has_minor_gridlines = False
        category_axis.has_major_gridlines = False
        category_axis.tick_labels.font.size = Pt(12)
        value_axis = chart.value_axis
        value_axis.has_minor_gridlines = False
        value_axis.has_major_gridlines = False
        value_axis.tick_labels.font.size = Pt(12)
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        if question.stat == 'percent':
            data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        data_labels.font.size = Pt(12)
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(12)
        chart.legend.include_in_layout = False

    def column_clustered_matrix_by_scale_points(self, question, slide):
        sub_questions = question.questions
        series_names = []  # responses
        series_lists = []  # frequencies for each scale point
        categories = []  # scale points
        is_first = True

        for sub_question in sub_questions:
            series_names.append(str(sub_question.prompt) + " (n=" + str(sub_question.n) + ")")
            if is_first:
                for response in sub_question.responses:
                    categories.append(response.response)
                question.stat = sub_question.stat
            is_first = False

        for name in series_names:
            frequencies = []
            series_lists.append([])
        response_idx = 0
        for sub_question in sub_questions:
            frequencies = series_lists[response_idx]
            for response in sub_question.responses:
                if response.has_frequency is True:
                    for year, frequency in response.frequencies.items():
                        frequencies.append(frequency)
                else:
                    frequencies.append(0)
            response_idx += 1

        chart_data = CategoryChartData()
        chart_data.categories = categories
        for idx in range(0, len(series_names)):
            chart_data.add_series(series_names[idx], iter(series_lists[idx]))

        left, top, width, height = Inches(6.75), Inches(2), Inches(6.5), Inches(5)
        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data).chart
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = "Q: " + str(question.prompt)
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)

        category_axis = chart.category_axis
        category_axis.has_minor_gridlines = False
        category_axis.has_major_gridlines = False
        category_axis.tick_labels.font.size = Pt(12)
        value_axis = chart.value_axis
        value_axis.has_minor_gridlines = False
        value_axis.has_major_gridlines = False
        value_axis.tick_labels.font.size = Pt(12)
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        if question.stat == 'percent':
            data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        data_labels.font.size = Pt(12)
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(12)
        chart.legend.include_in_layout = False

    def bar_clustered_matrix_by_categories(self, question, slide):
        sub_questions = question.questions
        series_names = []  # scale points
        series_lists = []  # frequencies for each scale point
        categories = []  # responses for matrix
        is_first = True

        for sub_question in sub_questions:
            categories.append(str(sub_question.prompt) + " (n=" + str(sub_question.n) + ")")
            if is_first:
                for response in sub_question.responses:
                    series_names.append(response.response)
                question.stat = sub_question.stat
            is_first = False

        for name in series_names:
            frequencies = []
            series_lists.append([])

        for sub_question in sub_questions:
            response_idx = 0
            for response in sub_question.responses:
                frequencies = series_lists[response_idx]
                if response.has_frequency is True:
                    for year, frequency in response.frequencies.items():
                        frequencies.append(frequency)
                else:
                    frequencies.append(0)
                response_idx += 1

        chart_data = CategoryChartData()
        chart_data.categories = categories
        for idx in range(0, len(series_names)):
            chart_data.add_series(series_names[idx], iter(series_lists[idx]))
        left, top, width, height = Inches(0), Inches(2), Inches(6.5), Inches(5)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data).chart
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = "Q: " + str(question.prompt)
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)

        category_axis = chart.category_axis
        category_axis.has_minor_gridlines = False
        category_axis.has_major_gridlines = False
        category_axis.tick_labels.font.size = Pt(12)
        value_axis = chart.value_axis
        value_axis.has_minor_gridlines = False
        value_axis.has_major_gridlines = False
        value_axis.tick_labels.font.size = Pt(12)
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        if question.stat == 'percent':
            data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        data_labels.font.size = Pt(12)
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(12)
        chart.legend.include_in_layout = False

    def bar_clustered_matrix_by_scale_points(self, question, slide):
        sub_questions = question.questions
        series_names = []  # responses
        series_lists = []  # frequencies for each scale point
        categories = []  # scale points
        is_first = True

        for sub_question in sub_questions:
            series_names.append(str(sub_question.prompt) + " (n=" + str(sub_question.n) + ")")
            if is_first:
                for response in sub_question.responses:
                    categories.append(response.response)
                question.stat = sub_question.stat
            is_first = False

        for name in series_names:
            frequencies = []
            series_lists.append([])
        response_idx = 0
        for sub_question in sub_questions:
            frequencies = series_lists[response_idx]
            for response in sub_question.responses:
                if response.has_frequency is True:
                    for year, frequency in response.frequencies.items():
                        frequencies.append(frequency)
                else:
                    frequencies.append(0)
            response_idx += 1

        chart_data = CategoryChartData()
        chart_data.categories = categories
        for idx in range(0, len(series_names)):
            chart_data.add_series(series_names[idx], iter(series_lists[idx]))
        left, top, width, height = Inches(6.75), Inches(2), Inches(6.5), Inches(5)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data).chart
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = "Q: " + str(question.prompt)
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)

        category_axis = chart.category_axis
        category_axis.has_minor_gridlines = False
        category_axis.has_major_gridlines = False
        category_axis.tick_labels.font.size = Pt(12)
        value_axis = chart.value_axis
        value_axis.has_minor_gridlines = False
        value_axis.has_major_gridlines = False
        value_axis.tick_labels.font.size = Pt(12)
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        if question.stat == 'percent':
            data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        data_labels.font.size = Pt(12)
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(12)
        chart.legend.include_in_layout = False

    def stacked_bar_matrix(self, question, slide, forward, diverge, side):
        sub_questions = question.questions
        series_names = []  # scale points
        series_lists = []  # frequencies for each scale point
        categories = []  # responses for matrix
        is_first = True

        for sub_question in sub_questions:
            categories.append(str(sub_question.prompt) + " (n=" + str(sub_question.n) + ")")
            if is_first:
                for response in sub_question.responses:
                    series_names.append(response.response)
                question.stat = sub_question.stat
            is_first = False

        for name in series_names:
            frequencies = []
            series_lists.append([])

        for sub_question in sub_questions:
            response_idx = 0
            for response in sub_question.responses:
                frequencies = series_lists[response_idx]
                if response.has_frequency is True:
                    for year, frequency in response.frequencies.items():
                        frequencies.append(frequency)
                else:
                    frequencies.append(0)
                response_idx += 1

        chart_data = CategoryChartData()
        chart_data.categories = categories
        for idx in range(0, len(series_names)):
            chart_data.add_series(series_names[idx], iter(series_lists[idx]))

        if side == 'top':
            left, top, width, height = Inches(1.5), Inches(1), Inches(10), Inches(3)
        if side == 'bottom':
            left, top, width, height = Inches(1.5), Inches(4), Inches(10), Inches(3)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_STACKED, left, top, width, height, chart_data).chart
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = "Q: " + str(question.prompt)
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
        chart.has_legend = True
        category_axis = chart.category_axis
        category_axis.has_minor_gridlines = False
        category_axis.has_major_gridlines = False
        category_axis.tick_labels.font.size = Pt(12)
        value_axis = chart.value_axis
        value_axis.has_minor_gridlines = False
        value_axis.has_major_gridlines = False
        value_axis.visible = False

        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(12)
        chart.legend.include_in_layout = False
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        if question.stat == 'percent':
            data_labels.number_format = '0%'
        data_labels.font.size = Pt(12)
        plot = chart.plots[0]

        length = len(plot.series)
        self.determine_color_indicies(length, forward, diverge)
        color_index = 0
        for series in plot.series:
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = self.colors_to_use[color_index]
            if fill.fore_color.rgb == RGBColor(0, 35, 100) or fill.fore_color.rgb == RGBColor(130, 130, 140):
                # make text white
                series.data_labels.font.color.rgb = RGBColor(255, 255, 255)
                series.data_labels.font.size =  Pt(12)
                series.data_labels.number_format = '0%'
                series.data_labels.show_value = True
            color_index += 1


    def column_clustered_allocation(self, question, slide):
        sub_questions = question.questions
        categories = []
        frequencies = []

        for sub_question in sub_questions:
            for response in sub_question.responses:
                categories.append(response.response + " (n=" + str(sub_question.n) + ")")
                for year, frequency in response.frequencies.items():
                    frequencies.append(frequency)

        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series('Series 1', iter(frequencies))
        left, top, width, height = Inches(0), Inches(2), Inches(6.5), Inches(5)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data).chart
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = "Q: " + str(question.prompt)
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
        chart.has_legend = False
        category_axis = chart.category_axis
        category_axis.has_minor_gridlines = False
        category_axis.has_major_gridlines = False
        category_axis.tick_labels.font.size = Pt(12)
        value_axis = chart.value_axis
        value_axis.has_minor_gridlines = False
        value_axis.has_major_gridlines = False
        value_axis.tick_labels.font.size = Pt(12)
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        data_labels.font.size = Pt(12)
        plot = chart.plots[0]
        series = plot.series[0]
        fill = series.format.fill
        fill.solid()
        fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1

    def bar_clustered_allocation(self, question, slide):
        sub_questions = question.questions
        categories = []
        frequencies = []

        for sub_question in sub_questions:
            for response in sub_question.responses:
                categories.append(response.response + " (n=" + str(sub_question.n) + ")")
                for year, frequency in response.frequencies.items():
                    frequencies.append(frequency)

        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series('Series 1', iter(frequencies))
        left, top, width, height = Inches(6.75), Inches(2), Inches(6.5), Inches(5)

        chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data).chart
        chart.chart_title.has_text_frame = True
        chart.chart_title.text_frame.text = "Q: " + str(question.prompt)
        chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
        chart.has_legend = False
        category_axis = chart.category_axis
        category_axis.has_minor_gridlines = False
        category_axis.has_major_gridlines = False
        category_axis.tick_labels.font.size = Pt(12)
        value_axis = chart.value_axis
        value_axis.has_minor_gridlines = False
        value_axis.has_major_gridlines = False
        value_axis.tick_labels.font.size = Pt(12)
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        # data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        data_labels.font.size = Pt(12)
        plot = chart.plots[0]
        series = plot.series[0]
        fill = series.format.fill
        fill.solid()
        fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1

    def determine_color_indicies(self, length, forward, diverge):
        self.colors_to_use = []
        if forward and diverge:
            if length == 1:
                self.colors_to_use.append(self.diverge_color_array[0])
            elif length == 2:
                self.colors_to_use.append(self.diverge_color_array[0])
                self.colors_to_use.append(self.diverge_color_array[9])
            elif length == 3:
                self.colors_to_use.append(self.diverge_color_array[0])
                self.colors_to_use.append(self.diverge_color_array[4])
                self.colors_to_use.append(self.diverge_color_array[9])
            elif length == 4:
                self.colors_to_use.append(self.diverge_color_array[0])
                self.colors_to_use.append(self.diverge_color_array[3])
                self.colors_to_use.append(self.diverge_color_array[4])
                self.colors_to_use.append(self.diverge_color_array[9])
            elif length == 5:
                self.colors_to_use.append(self.diverge_color_array[0])
                self.colors_to_use.append(self.diverge_color_array[3])
                self.colors_to_use.append(self.diverge_color_array[4])
                self.colors_to_use.append(self.diverge_color_array[5])
                self.colors_to_use.append(self.diverge_color_array[9])
            elif length == 6:
                self.colors_to_use.append(self.diverge_color_array[0])
                self.colors_to_use.append(self.diverge_color_array[1])
                self.colors_to_use.append(self.diverge_color_array[3])
                self.colors_to_use.append(self.diverge_color_array[4])
                self.colors_to_use.append(self.diverge_color_array[5])
                self.colors_to_use.append(self.diverge_color_array[9])
            elif length == 7:
                self.colors_to_use.append(self.diverge_color_array[0])
                self.colors_to_use.append(self.diverge_color_array[1])
                self.colors_to_use.append(self.diverge_color_array[3])
                self.colors_to_use.append(self.diverge_color_array[4])
                self.colors_to_use.append(self.diverge_color_array[5])
                self.colors_to_use.append(self.diverge_color_array[8])
                self.colors_to_use.append(self.diverge_color_array[9])
            elif length == 8:
                self.colors_to_use.append(self.diverge_color_array[0])
                self.colors_to_use.append(self.diverge_color_array[1])
                self.colors_to_use.append(self.diverge_color_array[3])
                self.colors_to_use.append(self.diverge_color_array[4])
                self.colors_to_use.append(self.diverge_color_array[5])
                self.colors_to_use.append(self.diverge_color_array[7])
                self.colors_to_use.append(self.diverge_color_array[8])
                self.colors_to_use.append(self.diverge_color_array[9])
            elif length == 9:
                self.colors_to_use.append(self.diverge_color_array[0])
                self.colors_to_use.append(self.diverge_color_array[1])
                self.colors_to_use.append(self.diverge_color_array[2])
                self.colors_to_use.append(self.diverge_color_array[3])
                self.colors_to_use.append(self.diverge_color_array[4])
                self.colors_to_use.append(self.diverge_color_array[5])
                self.colors_to_use.append(self.diverge_color_array[7])
                self.colors_to_use.append(self.diverge_color_array[8])
                self.colors_to_use.append(self.diverge_color_array[9])
            elif length ==10:
                self.colors_to_use.append(self.diverge_color_array[0])
                self.colors_to_use.append(self.diverge_color_array[1])
                self.colors_to_use.append(self.diverge_color_array[2])
                self.colors_to_use.append(self.diverge_color_array[3])
                self.colors_to_use.append(self.diverge_color_array[4])
                self.colors_to_use.append(self.diverge_color_array[5])
                self.colors_to_use.append(self.diverge_color_array[6])
                self.colors_to_use.append(self.diverge_color_array[7])
                self.colors_to_use.append(self.diverge_color_array[8])
                self.colors_to_use.append(self.diverge_color_array[9])

        if not forward and diverge:
            if length == 1:
                self.colors_to_use.append(self.diverge_color_array[0])
            elif length == 2:
                self.colors_to_use.append(self.diverge_color_array[9])
                self.colors_to_use.append(self.diverge_color_array[0])
            elif length == 3:
                self.colors_to_use.append(self.diverge_color_array[9])
                self.colors_to_use.append(self.diverge_color_array[4])
                self.colors_to_use.append(self.diverge_color_array[0])
            elif length == 4:
                self.colors_to_use.append(self.diverge_color_array[9])
                self.colors_to_use.append(self.diverge_color_array[4])
                self.colors_to_use.append(self.diverge_color_array[3])
                self.colors_to_use.append(self.diverge_color_array[0])
            elif length == 5:
                self.colors_to_use.append(self.diverge_color_array[9])
                self.colors_to_use.append(self.diverge_color_array[5])
                self.colors_to_use.append(self.diverge_color_array[4])
                self.colors_to_use.append(self.diverge_color_array[3])
                self.colors_to_use.append(self.diverge_color_array[0])
            elif length == 6:
                self.colors_to_use.append(self.diverge_color_array[9])
                self.colors_to_use.append(self.diverge_color_array[5])
                self.colors_to_use.append(self.diverge_color_array[4])
                self.colors_to_use.append(self.diverge_color_array[3])
                self.colors_to_use.append(self.diverge_color_array[1])
                self.colors_to_use.append(self.diverge_color_array[0])
            elif length == 7:
                self.colors_to_use.append(self.diverge_color_array[9])
                self.colors_to_use.append(self.diverge_color_array[8])
                self.colors_to_use.append(self.diverge_color_array[5])
                self.colors_to_use.append(self.diverge_color_array[4])
                self.colors_to_use.append(self.diverge_color_array[3])
                self.colors_to_use.append(self.diverge_color_array[1])
                self.colors_to_use.append(self.diverge_color_array[0])
            elif length == 8:
                self.colors_to_use.append(self.diverge_color_array[9])
                self.colors_to_use.append(self.diverge_color_array[8])
                self.colors_to_use.append(self.diverge_color_array[7])
                self.colors_to_use.append(self.diverge_color_array[5])
                self.colors_to_use.append(self.diverge_color_array[4])
                self.colors_to_use.append(self.diverge_color_array[3])
                self.colors_to_use.append(self.diverge_color_array[1])
                self.colors_to_use.append(self.diverge_color_array[0])
            elif length == 9:
                self.colors_to_use.append(self.diverge_color_array[9])
                self.colors_to_use.append(self.diverge_color_array[8])
                self.colors_to_use.append(self.diverge_color_array[7])
                self.colors_to_use.append(self.diverge_color_array[5])
                self.colors_to_use.append(self.diverge_color_array[4])
                self.colors_to_use.append(self.diverge_color_array[3])
                self.colors_to_use.append(self.diverge_color_array[2])
                self.colors_to_use.append(self.diverge_color_array[1])
                self.colors_to_use.append(self.diverge_color_array[0])
            elif length == 10:
                self.colors_to_use.append(self.diverge_color_array[9])
                self.colors_to_use.append(self.diverge_color_array[8])
                self.colors_to_use.append(self.diverge_color_array[7])
                self.colors_to_use.append(self.diverge_color_array[6])
                self.colors_to_use.append(self.diverge_color_array[5])
                self.colors_to_use.append(self.diverge_color_array[4])
                self.colors_to_use.append(self.diverge_color_array[3])
                self.colors_to_use.append(self.diverge_color_array[2])
                self.colors_to_use.append(self.diverge_color_array[1])
                self.colors_to_use.append(self.diverge_color_array[0])
        if forward and not diverge:
            if length == 1:
                self.colors_to_use.append(self.mono_color_array[0])
            elif length == 2:
                self.colors_to_use.append(self.mono_color_array[0])
                self.colors_to_use.append(self.mono_color_array[4])
            elif length == 3:
                self.colors_to_use.append(self.mono_color_array[0])
                self.colors_to_use.append(self.mono_color_array[2])
                self.colors_to_use.append(self.mono_color_array[4])
            elif length == 4:
                self.colors_to_use.append(self.mono_color_array[0])
                self.colors_to_use.append(self.mono_color_array[2])
                self.colors_to_use.append(self.mono_color_array[4])
                self.colors_to_use.append(self.mono_color_array[5])
            elif length == 5:
                self.colors_to_use.append(self.mono_color_array[0])
                self.colors_to_use.append(self.mono_color_array[2])
                self.colors_to_use.append(self.mono_color_array[3])
                self.colors_to_use.append(self.mono_color_array[4])
                self.colors_to_use.append(self.mono_color_array[5])
            elif length == 6:
                self.colors_to_use.append(self.mono_color_array[0])
                self.colors_to_use.append(self.mono_color_array[1])
                self.colors_to_use.append(self.mono_color_array[2])
                self.colors_to_use.append(self.mono_color_array[3])
                self.colors_to_use.append(self.mono_color_array[4])
                self.colors_to_use.append(self.mono_color_array[5])
            elif length == 7:
                self.colors_to_use.append(self.mono_color_array[0])
                self.colors_to_use.append(self.mono_color_array[1])
                self.colors_to_use.append(self.mono_color_array[2])
                self.colors_to_use.append(self.mono_color_array[3])
                self.colors_to_use.append(self.mono_color_array[4])
                self.colors_to_use.append(self.mono_color_array[5])
                self.colors_to_use.append(self.mono_color_array[6])
            elif length == 8:
                self.colors_to_use.append(self.mono_color_array[0])
                self.colors_to_use.append(self.mono_color_array[1])
                self.colors_to_use.append(self.mono_color_array[2])
                self.colors_to_use.append(self.mono_color_array[3])
                self.colors_to_use.append(self.mono_color_array[4])
                self.colors_to_use.append(self.mono_color_array[5])
                self.colors_to_use.append(self.mono_color_array[6])
                self.colors_to_use.append(self.mono_color_array[7])
            elif length == 9:
                self.colors_to_use.append(self.mono_color_array[0])
                self.colors_to_use.append(self.mono_color_array[1])
                self.colors_to_use.append(self.mono_color_array[2])
                self.colors_to_use.append(self.mono_color_array[3])
                self.colors_to_use.append(self.mono_color_array[4])
                self.colors_to_use.append(self.mono_color_array[5])
                self.colors_to_use.append(self.mono_color_array[6])
                self.colors_to_use.append(self.mono_color_array[7])
                self.colors_to_use.append(self.mono_color_array[8])
            elif length == 10:
                self.colors_to_use.append(self.mono_color_array[0])
                self.colors_to_use.append(self.mono_color_array[1])
                self.colors_to_use.append(self.mono_color_array[2])
                self.colors_to_use.append(self.mono_color_array[3])
                self.colors_to_use.append(self.mono_color_array[4])
                self.colors_to_use.append(self.mono_color_array[5])
                self.colors_to_use.append(self.mono_color_array[6])
                self.colors_to_use.append(self.mono_color_array[7])
                self.colors_to_use.append(self.mono_color_array[8])
                self.colors_to_use.append(self.mono_color_array[9])

        if not forward and not diverge:
            if length == 1:
                self.colors_to_use.append(self.mono_color_array[0])
            elif length == 2:
                self.colors_to_use.append(self.mono_color_array[4])
                self.colors_to_use.append(self.mono_color_array[0])
            elif length == 3:
                self.colors_to_use.append(self.mono_color_array[4])
                self.colors_to_use.append(self.mono_color_array[2])
                self.colors_to_use.append(self.mono_color_array[0])
            elif length == 4:
                self.colors_to_use.append(self.mono_color_array[5])
                self.colors_to_use.append(self.mono_color_array[4])
                self.colors_to_use.append(self.mono_color_array[2])
                self.colors_to_use.append(self.mono_color_array[0])
            elif length == 5:
                self.colors_to_use.append(self.mono_color_array[5])
                self.colors_to_use.append(self.mono_color_array[4])
                self.colors_to_use.append(self.mono_color_array[3])
                self.colors_to_use.append(self.mono_color_array[2])
                self.colors_to_use.append(self.mono_color_array[0])
            elif length == 6:
                self.colors_to_use.append(self.mono_color_array[5])
                self.colors_to_use.append(self.mono_color_array[4])
                self.colors_to_use.append(self.mono_color_array[3])
                self.colors_to_use.append(self.mono_color_array[2])
                self.colors_to_use.append(self.mono_color_array[1])
                self.colors_to_use.append(self.mono_color_array[0])
            elif length == 7:
                self.colors_to_use.append(self.mono_color_array[6])
                self.colors_to_use.append(self.mono_color_array[5])
                self.colors_to_use.append(self.mono_color_array[4])
                self.colors_to_use.append(self.mono_color_array[3])
                self.colors_to_use.append(self.mono_color_array[2])
                self.colors_to_use.append(self.mono_color_array[1])
                self.colors_to_use.append(self.mono_color_array[0])
            elif length == 8:
                self.colors_to_use.append(self.mono_color_array[7])
                self.colors_to_use.append(self.mono_color_array[6])
                self.colors_to_use.append(self.mono_color_array[5])
                self.colors_to_use.append(self.mono_color_array[4])
                self.colors_to_use.append(self.mono_color_array[3])
                self.colors_to_use.append(self.mono_color_array[2])
                self.colors_to_use.append(self.mono_color_array[1])
                self.colors_to_use.append(self.mono_color_array[0])
            elif length == 9:
                self.colors_to_use.append(self.mono_color_array[8])
                self.colors_to_use.append(self.mono_color_array[7])
                self.colors_to_use.append(self.mono_color_array[6])
                self.colors_to_use.append(self.mono_color_array[5])
                self.colors_to_use.append(self.mono_color_array[4])
                self.colors_to_use.append(self.mono_color_array[3])
                self.colors_to_use.append(self.mono_color_array[2])
                self.colors_to_use.append(self.mono_color_array[1])
                self.colors_to_use.append(self.mono_color_array[0])
            elif length == 10:
                self.colors_to_use.append(self.mono_color_array[9])
                self.colors_to_use.append(self.mono_color_array[8])
                self.colors_to_use.append(self.mono_color_array[7])
                self.colors_to_use.append(self.mono_color_array[6])
                self.colors_to_use.append(self.mono_color_array[5])
                self.colors_to_use.append(self.mono_color_array[4])
                self.colors_to_use.append(self.mono_color_array[3])
                self.colors_to_use.append(self.mono_color_array[2])
                self.colors_to_use.append(self.mono_color_array[1])
                self.colors_to_use.append(self.mono_color_array[0])





