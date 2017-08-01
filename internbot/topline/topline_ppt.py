from pptx import Presentation
from pptx.util import Inches, Pt

class ToplinePPT(object):
    
    def __init__ (self, questions, path_to_template):
        self.ppt = Presentation(path_to_template)
        self.questions = questions

    def save(self, path_to_output):
        self.write_questions()
        self.save_file(path_to_output)

    def write_questions(self):
        for question in self.questions:
            if question.type == 'Composite':
                self.write_composite_slide(question)
            else:
                self.write_basic_slide(question)

    def save_file(self, path_to_output):
        self.ppt.save(path_to_output)

    def write_composite_slide(self, question):
        slide = self.ppt.slides.add_slide(self.ppt.slide_layouts[3])
        self.write_name(question.name, slide)
        self.write_prompts(question, slide)

    def write_basic_slide(self, question):
        slide = self.ppt.slides.add_slide(self.ppt.slide_layouts[1])
        self.write_name(question.name, slide)
        self.write_prompt(question.prompt, slide)

    def write_name(self, name, slide):
        title = slide.shapes.title
        title.text = name

    def write_prompt(self, prompt, slide):
        body = slide.placeholders[1]
        body.text = prompt

    def write_prompts(self, composite_question, slide):
        overall_prompt = slide.placeholders[1]
        overall_prompt.text = composite_question.prompt
        sub_prompt = slide.placeholders[2]
        sub_prompt.text = ''
        for sub_question in composite_question.questions:
            sub_prompt.text += sub_question.prompt + '\n'

    def get_ppt(self):
        return self.ppt
        